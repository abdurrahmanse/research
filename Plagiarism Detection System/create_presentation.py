from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme
BLUE = RGBColor(52, 152, 219)
RED = RGBColor(231, 76, 60)
GREEN = RGBColor(46, 204, 113)
DARK_GRAY = RGBColor(44, 62, 80)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLUE
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(60)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(2))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content_dict):
    """Add a content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # White background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)
    
    # Title bar
    title_shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.9))
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = BLUE
    title_shape.line.color.rgb = BLUE
    
    # Title text
    title_frame = title_shape.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_frame.margin_left = Inches(0.3)
    
    # Content
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(6))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    for key, value in content_dict.items():
        p = text_frame.add_paragraph()
        p.text = key
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY
        p.level = 0
        p.space_before = Pt(12)
        
        if isinstance(value, list):
            for item in value:
                p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(13)
                p.font.color.rgb = DARK_GRAY
                p.level = 1
                p.space_before = Pt(4)
        else:
            p = text_frame.add_paragraph()
            p.text = str(value)
            p.font.size = Pt(13)
            p.font.color.rgb = DARK_GRAY
            p.level = 1

# Slide 1: Title
add_title_slide(prs, 
    "Plagiarism Detection System",
    "A Comprehensive Machine Learning Analysis\nComparing TF-IDF, BiLSTM, and BERT Models")

# Slide 2: Overview
add_content_slide(prs, "Project Overview", {
    "Objective": [
        "Build a comprehensive plagiarism detection system",
        "Compare three different ML/DL approaches",
        "Analyze performance and trade-offs"
    ],
    "Dataset": [
        "100 unique academic paper summaries",
        "197 text pairs with plagiarism labels",
        "Balanced: 99 plagiarized, 98 non-plagiarized"
    ],
    "Data Split": [
        "Training: 70% | Validation: 15% | Testing: 15%"
    ]
})

# Slide 3: Dataset
add_content_slide(prs, "Dataset Analysis", {
    "Text Statistics": [
        "Average length: 152 words",
        "Range: 45 - 251 words",
        "Perfectly balanced classes"
    ],
    "Preprocessing": [
        "Lowercased all text",
        "Removed URLs and emails",
        "Removed special characters",
        "Normalized whitespace"
    ],
    "Quality": [
        "No missing values",
        "Stratified splitting",
        "Ready for training"
    ]
})

# Slide 4: TF-IDF
add_content_slide(prs, "Model 1: TF-IDF (Baseline)", {
    "What it Does": [
        "Converts text to word frequency vectors",
        "Compares similarity using cosine distance",
        "Simple, fast, and interpretable"
    ],
    "Performance (Test Set)": [
        "Accuracy: 76.67% | Precision: 75.00%",
        "Recall: 77.78% | F1-Score: 76.36%",
        "ROC-AUC: 78.96%"
    ],
    "Pros": "Fast | Low memory | Good baseline"
})

# Slide 5: BiLSTM
add_content_slide(prs, "Model 2: BiLSTM (Deep Learning)", {
    "What it Does": [
        "Creates word embeddings (learns meaning)",
        "Uses bidirectional LSTM layers",
        "Understands context and sequences"
    ],
    "Performance (Test Set)": [
        "Accuracy: 83.33% | Precision: 84.21%",
        "Recall: 83.33% | F1-Score: 83.76%",
        "ROC-AUC: 89.41%"
    ],
    "Pros": "Better accuracy | Good balance | Context aware"
})

# Slide 6: BERT
add_content_slide(prs, "Model 3: BERT (State-of-the-Art)", {
    "What it Does": [
        "Pre-trained transformer architecture",
        "Bidirectional attention mechanism",
        "Understands deep semantic meaning"
    ],
    "Performance (Test Set)": [
        "Accuracy: 96.67% | Precision: 94.44%",
        "Recall: 100.00% ⭐ | F1-Score: 97.14%",
        "ROC-AUC: 99.51%"
    ],
    "Best For": "Maximum accuracy | Catches ALL plagiarism"
})

# Slide 7: Comparison
add_content_slide(prs, "Performance Comparison", {
    "Accuracy Ranking": [
        "🥇 BERT: 96.67%",
        "🥈 BiLSTM: 83.33%",
        "🥉 TF-IDF: 76.67%"
    ],
    "Recall Ranking": [
        "🥇 BERT: 100% (Catches ALL plagiarism)",
        "🥈 BiLSTM: 83.33%",
        "🥉 TF-IDF: 77.78%"
    ],
    "Speed Ranking": [
        "🥇 TF-IDF: ~10ms | 🥈 BiLSTM: ~100ms",
        "🥉 BERT: 1-2 seconds"
    ]
})

# Slide 8: Key Findings
add_content_slide(prs, "Key Findings & Insights", {
    "Model Effectiveness": [
        "BERT significantly outperforms both models",
        "BiLSTM provides excellent balance",
        "TF-IDF works well as baseline"
    ],
    "Recall is Critical": [
        "BERT's 100% recall = NO plagiarism missed",
        "94.44% precision = Most alerts are real",
        "99.51% ROC-AUC = Excellent discrimination"
    ],
    "Trade-offs": [
        "Accuracy vs Speed: Choose your priority",
        "Resources vs Performance: Budget matters"
    ]
})

# Slide 9: Recommendations
add_content_slide(prs, "Deployment Recommendations", {
    "For Academic Institutions": [
        "Recommended: BERT",
        "Reason: 100% recall (catches all plagiarism)",
        "Cost: GPU server or cloud deployment"
    ],
    "For Real-time Web Apps": [
        "Recommended: BiLSTM",
        "Reason: Fast enough, accurate enough",
        "Cost: Standard CPU server"
    ],
    "For Initial Screening": [
        "Recommended: TF-IDF Pipeline",
        "Reason: Very fast, minimal resources"
    ]
})

# Slide 10: Technical Insights
add_content_slide(prs, "Technical Learnings", {
    "Text Processing": [
        "Data cleaning directly impacts performance",
        "Standardization ensures consistent input",
        "Successfully handled 200+ documents"
    ],
    "Model Architecture": [
        "Embeddings capture semantic meaning",
        "BiLSTM effective for sequences",
        "Transfer learning (BERT) saves time"
    ],
    "Evaluation": [
        "Multiple metrics needed for full picture",
        "Recall crucial for plagiarism detection",
        "Confusion matrices reveal patterns"
    ]
})

# Slide 11: Challenges
add_content_slide(prs, "Challenges & Solutions", {
    "Challenge 1: BERT Tokenizer Error": [
        "Problem: BertTokenizer not callable",
        "Solution: Switched to AutoTokenizer",
        "Result: Stable, compatible tokenization"
    ],
    "Challenge 2: Data Balance": [
        "Problem: Potential class imbalance",
        "Solution: Stratified splitting",
        "Result: Balanced across all sets"
    ],
    "Challenge 3: Computation": [
        "Problem: BERT training slow",
        "Solution: Optimized batch size",
        "Result: Reasonable time"
    ]
})

# Slide 12: Future Work
add_content_slide(prs, "Future Enhancements", {
    "Short-term": [
        "Cross-validation for robustness",
        "Hyperparameter tuning",
        "Ensemble method (combine all 3)",
        "Data augmentation"
    ],
    "Medium-term": [
        "Production API deployment",
        "Real-time web application",
        "Multilingual support",
        "Domain-specific fine-tuning"
    ],
    "Long-term": [
        "LLM integration (GPT-4)",
        "Explainable AI features",
        "Active learning pipeline",
        "Few-shot learning"
    ]
})

# Slide 13: Conclusion
add_content_slide(prs, "Conclusions", {
    "Project Success": [
        "✓ Built 3 working models successfully",
        "✓ BERT achieves 96.67% accuracy",
        "✓ Comprehensive analysis completed",
        "✓ Actionable recommendations provided"
    ],
    "Key Takeaways": [
        "Trade-offs between accuracy & speed exist",
        "Proper evaluation metrics are critical",
        "BERT excels; BiLSTM offers balance",
        "Context matters in plagiarism detection"
    ],
    "Winner": [
        "🏆 BERT: 96.67% accuracy, 100% recall"
    ]
})

# Slide 14: Q&A
slide = prs.slides.add_slide(prs.slide_layouts[6])
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = DARK_GRAY

title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1.5))
title_frame = title_box.text_frame
title_frame.text = "Questions & Discussion"
title_frame.paragraphs[0].font.size = Pt(50)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.8), Inches(9), Inches(1.5))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Plagiarism Detection System:\nMachine Learning Analysis"
subtitle_frame.paragraphs[0].font.size = Pt(20)
subtitle_frame.paragraphs[0].font.color.rgb = BLUE
subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save
output_path = '/Users/abdurrahman/Downloads/test/Plagiarism_Detection_Presentation.pptx'
prs.save(output_path)

print(f"✅ PowerPoint presentation created successfully!")
print(f"📁 File: Plagiarism_Detection_Presentation.pptx")
print(f"📊 Total slides: {len(prs.slides)}")
print("\n📋 Slides included:")
slides = [
    "1. Title Slide", "2. Project Overview", "3. Dataset Analysis",
    "4. TF-IDF Model", "5. BiLSTM Model", "6. BERT Model",
    "7. Performance Comparison", "8. Key Findings",
    "9. Deployment Recommendations", "10. Technical Learnings",
    "11. Challenges & Solutions", "12. Future Enhancements",
    "13. Conclusions", "14. Q&A Slide"
]
for slide in slides:
    print(f"   {slide}")
